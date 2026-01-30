# backend/rag_engine.py - Moteur RAG Complet
"""
Module de recherche RAG (Retrieval-Augmented Generation)
- Indexation automatique documents
- Recherche sémantique FAISS
- Génération réponses avec Claude (free: Haiku, payant: Sonnet)
"""

import os
import json
import time
from pathlib import Path
from typing import List, Dict, Any, Optional
import numpy as np
import faiss
from openai import OpenAI
from datetime import datetime
from dotenv import load_dotenv

# NOUVEAU : Import Claude
import anthropic

# Extraction documents
import PyPDF2
from PIL import Image
import pytesseract
from pptx import Presentation
import docx

# ============================================
# CONFIGURATION
# ============================================

# Charger les variables d'environnement depuis .env
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    raise ValueError("❌ OPENAI_API_KEY non trouvée ! Créez un fichier .env avec votre clé API.")

ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
if not ANTHROPIC_API_KEY:
    raise ValueError("❌ ANTHROPIC_API_KEY non trouvée ! Créez un fichier .env avec votre clé API.")

# Clients
openai_client = OpenAI(api_key=OPENAI_API_KEY)
claude_client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# Garder OpenAI pour embeddings (meilleurs pour recherche sémantique)
EMBEDDING_MODEL = "text-embedding-3-large"

# NOUVEAU : Modèles Claude selon le plan
PLAN_MODELS = {
    "free": "claude-3-haiku-20240307",       # Haiku pour free
    "basic": "claude-sonnet-4-20250514",     # Sonnet 4 pour payants
    "pro": "claude-sonnet-4-20250514",
    "premium": "claude-sonnet-4-20250514"
}

# ============================================
# SYSTEM PROMPT BILINGUE - Anti-hallucination STRICT
# ============================================

def get_bilingual_system_prompt(topic: str, language: str = "fr") -> str:
    """
    Génère le system prompt bilingue avec le sujet du cours
    
    Args:
        topic: Sujet du cours
        language: Langue de sortie ("fr" ou "en")
    """
    
    if language == "en":
        return f"""You are a bilingual study assistant for StudyGenie, working on the course "{topic}".

═══════════════════════════════════════════════════════════════
LANGUAGE RULE (ABSOLUTE PRIORITY - NEVER VIOLATE)
═══════════════════════════════════════════════════════════════
• The user has chosen ENGLISH as their language.
• You MUST respond ONLY and ENTIRELY in ENGLISH.
• If the course content (CONTEXT) is in French or another language:
  1. Understand the content in its original language
  2. Translate concepts FAITHFULLY to English
  3. Respond ENTIRELY in English
• NEVER mix languages in your response.
• NEVER include French words (except proper nouns or untranslatable terms).

═══════════════════════════════════════════════════════════════
MATHEMATICAL FORMULAS & TECHNICAL TERMS (PRESERVE AS-IS)
═══════════════════════════════════════════════════════════════
• Keep mathematical notation EXACTLY as written:
  - lim, log, ln, sin, cos, tan, exp, etc.
  - Symbols: ∑, ∫, ∂, ∇, π, θ, α, β, γ, Δ, →, ∞
  - Formulas: f'(x), dy/dx, ∂f/∂x, etc.
• Keep standard abbreviations: CPU, API, SQL, HTML, etc.
• Keep proper nouns: Newton, Euler, Gauss, etc.

═══════════════════════════════════════════════════════════════
CONTENT RULES (STRICT - NEVER VIOLATE)
═══════════════════════════════════════════════════════════════
1. Base responses ONLY on the provided CONTEXT.
2. If information is NOT in CONTEXT: "No data available for this topic currently."
3. If question is outside course scope: "This question is outside the scope of '{topic}'."
4. Do NOT invent, extrapolate, or use general knowledge.
5. Do NOT generate placeholder content like "Concept A" or "Formula 1".
6. Cite page numbers when available: (Page X) or (Source: Page X)

═══════════════════════════════════════════════════════════════
OUTPUT FORMAT
═══════════════════════════════════════════════════════════════
• Language: 100% English (except formulas and proper nouns)
• Style: Clear, structured, pedagogical
• No unnecessary introductions or generic conclusions
• Faithful to the course content

You are evaluated on RELIABILITY and LANGUAGE CONSISTENCY."""

    else:  # French (default)
        return f"""Tu es un assistant d'étude bilingue pour StudyGenie, travaillant sur le cours "{topic}".

═══════════════════════════════════════════════════════════════
RÈGLE DE LANGUE (PRIORITÉ ABSOLUE - NE JAMAIS VIOLER)
═══════════════════════════════════════════════════════════════
• L'utilisateur a choisi le FRANÇAIS comme langue.
• Tu DOIS répondre UNIQUEMENT et ENTIÈREMENT en FRANÇAIS.
• Si le contenu du cours (CONTEXT) est en anglais ou autre langue:
  1. Comprendre le contenu dans sa langue originale
  2. Traduire FIDÈLEMENT les concepts en français
  3. Répondre ENTIÈREMENT en français
• Ne JAMAIS mélanger les langues dans ta réponse.
• Ne JAMAIS inclure de mots anglais (sauf noms propres ou termes intraduisibles).

═══════════════════════════════════════════════════════════════
FORMULES MATHÉMATIQUES & TERMES TECHNIQUES (PRÉSERVER TEL QUEL)
═══════════════════════════════════════════════════════════════
• Garder la notation mathématique EXACTEMENT comme écrite:
  - lim, log, ln, sin, cos, tan, exp, etc.
  - Symboles: ∑, ∫, ∂, ∇, π, θ, α, β, γ, Δ, →, ∞
  - Formules: f'(x), dy/dx, ∂f/∂x, etc.
• Garder les abréviations standard: CPU, API, SQL, HTML, etc.
• Garder les noms propres: Newton, Euler, Gauss, etc.

═══════════════════════════════════════════════════════════════
RÈGLES DE CONTENU (STRICT - NE JAMAIS VIOLER)
═══════════════════════════════════════════════════════════════
1. Baser les réponses UNIQUEMENT sur le CONTEXT fourni.
2. Si l'information n'est PAS dans le CONTEXT: "Aucune donnée disponible pour ce sujet actuellement."
3. Si la question est hors sujet: "Cette question est hors du sujet '{topic}'."
4. Ne PAS inventer, extrapoler ou utiliser des connaissances générales.
5. Ne PAS générer de contenu placeholder comme "Concept A" ou "Formule 1".
6. Citer les numéros de page quand disponibles: (Page X) ou (Source: Page X)

═══════════════════════════════════════════════════════════════
FORMAT DE SORTIE
═══════════════════════════════════════════════════════════════
• Langue: 100% Français (sauf formules et noms propres)
• Style: Clair, structuré, pédagogique
• Pas d'introductions inutiles ou conclusions génériques
• Fidèle au contenu du cours

Tu es évalué sur la FIABILITÉ et la COHÉRENCE LINGUISTIQUE."""


# Garder l'ancienne fonction pour compatibilité
def get_system_prompt(topic: str) -> str:
    """Ancienne fonction - redirige vers bilingue en français"""
    return get_bilingual_system_prompt(topic, "fr")


CHUNK_SIZE = 800
CHUNK_OVERLAP = 200


def get_claude_model(user_plan: str) -> str:
    """Retourne le modèle Claude selon le plan utilisateur"""
    return PLAN_MODELS.get(user_plan, PLAN_MODELS["free"])


# ============================================
# EXTRACTION TEXTE (VOTRE CODE - INCHANGÉ)
# ============================================

def extract_text_from_pdf(file_path: str) -> tuple[str, int]:
    """Extrait texte d'un PDF
    Returns: (text, page_count)
    """
    text = ""
    page_count = 0
    
    try:
        with open(file_path, 'rb') as f:
            pdf = PyPDF2.PdfReader(f)
            page_count = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                
                # Si pas de texte (PDF scanné), essayer OCR
                if not page_text.strip() and has_tesseract():
                    print(f"   Page {page_num+1} : PDF scanné, OCR...")
                    page_text = "[Page scannée - OCR à implémenter]"
                
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"
    
    except Exception as e:
        print(f"❌ Erreur extraction PDF : {e}")
        text = ""
    
    return text, page_count


def extract_text_from_image(file_path: str) -> str:
    """Extrait texte d'une image avec OCR"""
    try:
        if not has_tesseract():
            return "[OCR non disponible - Installer Tesseract]"
        
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='fra+eng')
        return text
    
    except Exception as e:
        print(f"❌ Erreur OCR image : {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> tuple[str, int]:
    """Extrait texte d'un PowerPoint
    Returns: (text, slide_count)
    """
    text = ""
    slide_count = 0
    
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            
            # Texte des shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            
            # Notes du présentateur
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    text += f"\n[Notes : {notes}]\n"
    
    except Exception as e:
        print(f"❌ Erreur extraction PPTX : {e}")
    
    return text, slide_count


def extract_text_from_docx(file_path: str) -> str:
    """Extrait texte d'un document Word"""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text
    
    except Exception as e:
        print(f"❌ Erreur extraction DOCX : {e}")
        return ""


def extract_text_from_txt(file_path: str) -> str:
    """Extrait texte d'un fichier .txt (ou transcription)"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"❌ Erreur lecture TXT : {e}")
        return ""


def has_tesseract() -> bool:
    """Vérifie si Tesseract est installé"""
    try:
        pytesseract.get_tesseract_version()
        return True
    except:
        return False


# ============================================
# TRAITEMENT TEXTE (VOTRE CODE - INCHANGÉ)
# ============================================

def clean_text(text: str) -> str:
    """Nettoie le texte"""
    import re
    
    # Enlever caractères spéciaux
    text = text.replace('\x00', '')
    text = text.replace('\r\n', '\n')
    
    # Espaces multiples → 1 espace
    text = re.sub(r' +', ' ', text)
    
    # Lignes vides multiples → 1 ligne vide
    text = re.sub(r'\n\n+', '\n\n', text)
    
    return text.strip()


def create_chunks(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[Dict]:
    """Découpe texte en chunks avec métadonnées"""
    
    chunks = []
    start = 0
    chunk_id = 0
    
    # Détection pages (si marqueurs présents)
    pages = text.split("--- Page ")
    
    if len(pages) > 1:
        # PDF avec pages marquées
        for page_num, page_content in enumerate(pages[1:], start=1):
            # Extraire numéro de page
            parts = page_content.split("---", 1)
            if len(parts) == 2:
                page_text = parts[1].strip()
            else:
                page_text = page_content.strip()
            
            # Découper cette page en chunks si nécessaire
            if len(page_text) <= chunk_size:
                chunks.append({
                    'id': chunk_id,
                    'content': page_text,
                    'page': page_num,
                    'length': len(page_text)
                })
                chunk_id += 1
            else:
                # Page trop longue, découper
                page_chunks = simple_chunk(page_text, chunk_size, overlap)
                for pc in page_chunks:
                    chunks.append({
                        'id': chunk_id,
                        'content': pc,
                        'page': page_num,
                        'length': len(pc)
                    })
                    chunk_id += 1
    else:
        # Slides ou pas de pages marquées
        slides = text.split("--- Slide ")
        
        if len(slides) > 1:
            for slide_num, slide_content in enumerate(slides[1:], start=1):
                parts = slide_content.split("---", 1)
                if len(parts) == 2:
                    slide_text = parts[1].strip()
                else:
                    slide_text = slide_content.strip()
                
                chunks.append({
                    'id': chunk_id,
                    'content': slide_text,
                    'slide': slide_num,
                    'length': len(slide_text)
                })
                chunk_id += 1
        else:
            # Découpage simple
            simple_chunks = simple_chunk(text, chunk_size, overlap)
            for i, chunk_text in enumerate(simple_chunks):
                chunks.append({
                    'id': i,
                    'content': chunk_text,
                    'length': len(chunk_text)
                })
    
    return chunks


def simple_chunk(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Découpage simple avec chevauchement"""
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        end = start + chunk_size
        chunk_text = text[start:end]
        chunks.append(chunk_text)
        start += (chunk_size - overlap)
    
    return chunks


# ============================================
# EMBEDDINGS (OPENAI - INCHANGÉ)
# ============================================

def get_embedding(text: str) -> List[float]:
    """Génère embedding d'un texte avec OpenAI (gardé pour qualité)"""
    try:
        response = openai_client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=text
        )
        return response.data[0].embedding
    
    except Exception as e:
        print(f"❌ Erreur embedding : {e}")
        return None


# ============================================
# INDEXATION (VOTRE CODE - INCHANGÉ)
# ============================================

def index_course(
    user_id: int,
    course_id: int,
    file_path: str,
    course_name: str
) -> Dict[str, Any]:
    """
    Index un cours complet
    VOTRE CODE ORIGINAL - INCHANGÉ
    """
    
    print(f"\n{'='*60}")
    print(f"📚 INDEXATION COURS : {course_name}")
    print(f"{'='*60}\n")
    
    start_time = time.time()
    
    # 1. Extraction texte
    print("📄 Extraction texte...")
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        text, page_count = extract_text_from_pdf(file_path)
    elif file_ext == '.pptx':
        text, page_count = extract_text_from_pptx(file_path)
    elif file_ext == '.docx':
        text = extract_text_from_docx(file_path)
        page_count = 1
    elif file_ext in ['.txt', '.md']:
        text = extract_text_from_txt(file_path)
        page_count = 1
    elif file_ext in ['.jpg', '.jpeg', '.png']:
        text = extract_text_from_image(file_path)
        page_count = 1
    else:
        raise ValueError(f"Format {file_ext} non supporté")
    
    if not text.strip():
        raise ValueError("❌ Aucun texte extrait du fichier")
    
    print(f"✅ Texte extrait : {len(text)} caractères, {page_count} pages")
    
    # 2. Nettoyage
    print("\n🧹 Nettoyage texte...")
    text = clean_text(text)
    
    # 3. Chunking
    print("\n✂️ Découpage en chunks...")
    chunks = create_chunks(text)
    print(f"✅ {len(chunks)} chunks créés")
    
    # 4. Embeddings
    print("\n🧠 Génération embeddings...")
    embeddings = []
    
    for i, chunk in enumerate(chunks):
        if (i + 1) % 10 == 0:
            print(f"   Progress: {i+1}/{len(chunks)}")
        
        embedding = get_embedding(chunk['content'])
        if embedding:
            embeddings.append(embedding)
        else:
            embeddings.append([0.0] * 3072)  # Embedding vide si erreur
    
    embeddings_np = np.array(embeddings, dtype='float32')
    
    # 5. Index FAISS
    print("\n🔍 Création index FAISS...")
    dimension = embeddings_np.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings_np)
    
    print(f"✅ Index créé : {index.ntotal} vecteurs")
    
    # 6. Sauvegarde
    print("\n💾 Sauvegarde index...")
    
    course_dir = Path(f"data/users/{user_id}/courses/{course_id}")
    course_dir.mkdir(parents=True, exist_ok=True)
    
    # Sauvegarder FAISS
    faiss.write_index(index, str(course_dir / "index.faiss"))
    
    # Sauvegarder chunks
    with open(course_dir / "chunks.json", 'w', encoding='utf-8') as f:
        json.dump(chunks, f, ensure_ascii=False, indent=2)
    
    # Métadonnées
    metadata = {
        'course_name': course_name,
        'indexed_at': datetime.now().isoformat(),
        'chunks_count': len(chunks),
        'page_count': page_count,
        'file_type': file_ext
    }
    
    with open(course_dir / "metadata.json", 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    elapsed_time = time.time() - start_time
    
    print(f"\n{'='*60}")
    print(f"✅ INDEXATION TERMINÉE EN {elapsed_time:.1f}s")
    print(f"{'='*60}\n")
    
    return {
        'success': True,
        'chunks_count': len(chunks),
        'page_count': page_count,
        'time': elapsed_time
    }


# ============================================
# RECHERCHE RAG (MODIFIÉ POUR CLAUDE)
# ============================================

def search_and_answer(
    user_id: int,
    course_id: int,
    question: str,
    user_plan: str = "free",
    top_k: int = 5,
    language: str = "fr"  # NOUVEAU: Langue de sortie
) -> Dict[str, Any]:
    """
    Recherche et génère une réponse avec Claude (BILINGUE)
    
    Args:
        user_id: ID utilisateur
        course_id: ID cours
        question: Question posée
        user_plan: Plan utilisateur ("free", "basic", "pro", "premium")
        top_k: Nombre de chunks à récupérer
        language: Langue de sortie ("fr" ou "en")
    """
    
    start_time = time.time()
    
    print(f"\n{'='*60}")
    print(f"❓ QUESTION : {question}")
    print(f"👤 PLAN : {user_plan}")
    print(f"🌐 LANGUE : {language.upper()}")
    print(f"{'='*60}")
    
    # Charger index et chunks
    course_dir = Path(f"data/users/{user_id}/courses/{course_id}")
    
    if not course_dir.exists():
        raise FileNotFoundError(f"Cours {course_id} non trouvé")
    
    # Charger métadonnées pour récupérer le nom du cours (TOPIC)
    course_topic = "Cours"  # Valeur par défaut
    metadata_path = course_dir / "metadata.json"
    if metadata_path.exists():
        with open(metadata_path, 'r', encoding='utf-8') as f:
            metadata = json.load(f)
            course_topic = metadata.get('course_name', 'Cours')
    
    print(f"📚 SUJET : {course_topic}")
    
    # Charger FAISS
    index = faiss.read_index(str(course_dir / "index.faiss"))
    
    # Charger chunks
    with open(course_dir / "chunks.json", 'r', encoding='utf-8') as f:
        chunks = json.load(f)
    
    print(f"✅ Index chargé : {index.ntotal} vecteurs")
    
    # Recherche chunks pertinents
    print("\n🔍 Recherche chunks pertinents...")
    
    response = openai_client.embeddings.create(
        model=EMBEDDING_MODEL,
        input=question
    )
    
    question_embedding = np.array([response.data[0].embedding], dtype='float32')
    distances, indices = index.search(question_embedding, top_k)
    
    # Récupérer chunks
    relevant_chunks = []
    
    for i, (dist, idx) in enumerate(zip(distances[0], indices[0])):
        chunk = chunks[idx]
        
        relevant_chunks.append({
            'text': chunk['content'],
            'page': chunk.get('page', chunk.get('slide', 'N/A')),
            'distance': float(dist),
            'rank': i + 1
        })
        
        print(f"   #{i+1} - Distance: {dist:.3f} - Page: {chunk.get('page', 'N/A')}")
    
    # Construire contexte
    context = "\n\n---\n\n".join([
        f"[Source - Page {c['page']}]\n{c['text']}"
        for c in relevant_chunks
    ])
    
    # ===== APPEL CLAUDE BILINGUE =====
    print(f"\n🤖 Génération réponse avec Claude {get_claude_model(user_plan)} en {language.upper()}...")
    
    # User message selon la langue
    if language == "en":
        user_message = f"""CONTEXT:
{context}

QUESTION:
{question}

STRICT ANSWER (based only on CONTEXT, in English):"""
    else:
        user_message = f"""CONTEXT :
{context}

QUESTION :
{question}

RÉPONSE STRICTE (basée uniquement sur le CONTEXT, en français) :"""
    
    model = get_claude_model(user_plan)
    
    claude_response = claude_client.messages.create(
        model=model,
        max_tokens=1200,
        temperature=0.1,
        system=get_bilingual_system_prompt(course_topic, language),
        messages=[
            {"role": "user", "content": user_message}
        ]
    )
    
    answer_text = claude_response.content[0].text
    tokens_used = claude_response.usage.input_tokens + claude_response.usage.output_tokens
    
    # Confiance
    avg_distance = float(np.mean(distances[0]))
    confidence_rag = 1.0 / (1.0 + avg_distance)
    
    response_time = time.time() - start_time
    
    print(f"\n{'='*60}")
    print(f"✅ RÉPONSE GÉNÉRÉE EN {response_time:.1f}s ({language.upper()})")
    print(f"🤖 Modèle utilisé : {model}")
    print(f"📊 Tokens utilisés : {tokens_used}")
    print(f"📊 Confiance : {confidence_rag:.2%}")
    print(f"{'='*60}\n")
    
    result = {
        'question': question,
        'answer': answer_text,
        'sources': [
            {
                'text': c['text'][:200] + '...' if len(c['text']) > 200 else c['text'],
                'page': c['page'],
                'confidence': float(1.0 / (1.0 + c['distance']))
            }
            for c in relevant_chunks[:3]
        ],
        'confidence': confidence_rag,
        'response_time': response_time,
        'model_used': model,
        'tokens_used': tokens_used,
        'user_plan': user_plan,
        'language': language
    }
    
    return result


# ============================================
# VERSION AMÉLIORÉE (AVEC VALIDATIONS)
# ============================================

# Vos fonctions de validation (inchangées)
def validate_answer_completeness(question: str, answer: str) -> Dict:
    """Valide si réponse est complète"""
    import re
    
    issues = []
    score = 100
    
    # Détecter questions multiples
    subquestions = re.findall(r'[a-z]\)', question.lower())
    
    if subquestions and len(subquestions) > 1:
        for sub in subquestions:
            if sub not in answer.lower():
                issues.append({
                    'type': 'missing_subquestion',
                    'message': f"Sous-question {sub} non traitée"
                })
                score -= 30
    
    # Phrase coupée
    if not answer.strip().endswith(('.', '!', '?', ':', ')')):
        issues.append({
            'type': 'incomplete_sentence',
            'message': "Réponse se termine abruptement"
        })
        score -= 20
    
    # Trop courte
    if len(answer.split()) < 20 and "n'est pas dans le cours" not in answer:
        issues.append({
            'type': 'too_short',
            'message': "Réponse trop courte"
        })
        score -= 10
    
    return {
        'is_complete': score >= 70,
        'score': max(0, score),
        'issues': issues
    }


def extract_methodologies_from_course(text: str) -> Dict:
    """Extrait méthodologies du cours"""
    import re
    
    methodologies = []
    formulas = []
    
    # Patterns méthodologies
    method_patterns = [
        r'(?:méthodologie|méthode|procédure|étapes?)\s*:?\s*\n((?:[-•\d]\s*.+\n)+)',
        r'(?:pour|afin de)\s+(?:calculer|résoudre|trouver)\s+.+?:\s*\n((?:[-•\d]\s*.+\n)+)'
    ]
    
    for pattern in method_patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE | re.MULTILINE)
        for match in matches:
            methodologies.append({
                'text': match.group(0).strip(),
                'position': match.start()
            })
    
    # Formules
    formula_patterns = [
        r'[A-Za-z_]\s*=\s*[^=\n]{5,50}',
        r'\\frac\{.+?\}\{.+?\}',
        r'[A-Za-z]\^[0-9\{]'
    ]
    
    for pattern in formula_patterns:
        matches = re.finditer(pattern, text)
        for match in matches:
            formulas.append({
                'text': match.group(0).strip()
            })
    
    return {
        'methodologies': methodologies[:5],
        'formulas': formulas[:10],
        'has_methodology': len(methodologies) > 0
    }


def detect_exercise_type(question: str, context: str) -> str:
    """Détecte le type d'exercice"""
    import re
    
    combined_text = (question + " " + context).lower()
    
    exercise_types = {
        'calculation': [r'calcul\w*', r'détermin\w*', r'\btrouv\w*\s+(?:la|le|l\')\s+\w+'],
        'demonstration': [r'démontr\w*', r'prouv\w*', r'justifi\w*'],
        'analysis': [r'analys\w*', r'compar\w*', r'expliqu\w*'],
        'application': [r'appliqu\w*', r'utilis\w*', r'cas\s+pratique']
    }
    
    for ex_type, patterns in exercise_types.items():
        for pattern in patterns:
            if re.search(pattern, combined_text, re.IGNORECASE):
                return ex_type
    
    return 'general'


def search_and_answer_improved(
    user_id: int,
    course_id: int,
    question: str,
    user_plan: str = "free",
    top_k: int = 5,
    language: str = "fr"  # NOUVEAU: Langue de sortie
) -> Dict[str, Any]:
    """
    Version améliorée avec validation DE NOTATIONS + BILINGUE
    MODIFIÉ : Utilise Claude avec support multilingue
    
    Args:
        user_id: ID utilisateur
        course_id: ID du cours
        question: Question de l'utilisateur
        user_plan: Plan (free/pro)
        top_k: Nombre de chunks à récupérer
        language: Langue de sortie ("fr" ou "en")
    """
    
    # Import du module de validation des notations
    from validation_notations import (
        extract_math_notations_from_context,
        validate_notation_consistency,
        build_notation_aware_prompt
    )
    
    start_time = time.time()
    
    print(f"\n{'='*60}")
    print(f"❓ QUESTION : {question}")
    print(f"👤 PLAN : {user_plan}")
    print(f"🌐 LANGUE : {language.upper()}")
    print(f"{'='*60}")
    
    # Charger index et chunks
    course_dir = Path(f"data/users/{user_id}/courses/{course_id}")
    
    if not course_dir.exists():
        raise FileNotFoundError(f"Cours {course_id} non trouvé")
    
    # Charger métadonnées pour récupérer le nom du cours (TOPIC)
    course_topic = "Cours"  # Valeur par défaut
    metadata_path = course_dir / "metadata.json"
    if metadata_path.exists():
        with open(metadata_path, 'r', encoding='utf-8') as f:
            metadata = json.load(f)
            course_topic = metadata.get('course_name', 'Cours')
    
    print(f"📚 SUJET : {course_topic}")
    
    # Charger FAISS
    index = faiss.read_index(str(course_dir / "index.faiss"))
    
    # Charger chunks
    with open(course_dir / "chunks.json", 'r', encoding='utf-8') as f:
        chunks = json.load(f)
    
    print(f"✅ Index chargé : {index.ntotal} vecteurs")
    
    # Recherche chunks pertinents
    print("\n🔍 Recherche chunks pertinents...")
    
    response = openai_client.embeddings.create(
        model=EMBEDDING_MODEL,
        input=question
    )
    
    question_embedding = np.array([response.data[0].embedding], dtype='float32')
    distances, indices = index.search(question_embedding, top_k)
    
    # Récupérer chunks
    relevant_chunks = []
    
    for i, (dist, idx) in enumerate(zip(distances[0], indices[0])):
        chunk = chunks[idx]
        
        relevant_chunks.append({
            'text': chunk['content'],
            'page': chunk.get('page', chunk.get('slide', 'N/A')),
            'distance': float(dist),
            'rank': i + 1
        })
        
        print(f"   #{i+1} - Distance: {dist:.3f} - Page: {chunk.get('page', 'N/A')}")
    
    # Analyser le contexte
    print("\n📚 Extraction méthodologie du cours...")
    
    context_text = "\n\n".join([c['text'] for c in relevant_chunks])
    
    methodologies_data = extract_methodologies_from_course(context_text)
    exercise_type = detect_exercise_type(question, context_text)
    
    # Extraire les notations mathématiques du cours
    print("\n📐 Extraction des notations mathématiques...")
    context_notations = extract_math_notations_from_context(context_text)
    notation_instructions = build_notation_aware_prompt(context_text, question)
    
    print(f"   Type d'exercice : {exercise_type}")
    print(f"   Méthodologies trouvées : {len(methodologies_data['methodologies'])}")
    print(f"   Notations détectées : {sum(len(v) for v in context_notations.values())} symboles")
    
    # Construire contexte
    context = "\n\n---\n\n".join([
        f"[Source - Page {c['page']}]\n{c['text']}"
        for c in relevant_chunks
    ])
    
    # Sections additionnelles
    methodology_section = ""
    if methodologies_data['methodologies']:
        methodology_section = "\n\n🔹 MÉTHODOLOGIE DU COURS :\n"
        for method in methodologies_data['methodologies'][:2]:
            methodology_section += f"{method['text']}\n"
    
    formulas_section = ""
    if methodologies_data['formulas']:
        formulas_section = "\n\n🔹 FORMULES DU COURS :\n"
        for formula in methodologies_data['formulas'][:5]:
            formulas_section += f"• {formula['text']}\n"
    
    # Prompt amélioré - utilise le SYSTEM PROMPT BILINGUE avec TOPIC et LANGUE
    system_prompt = get_bilingual_system_prompt(course_topic, language)
    
    # Ajouter règles additionnelles selon la langue
    if language == "en":
        system_prompt += """

🎯 ADDITIONAL RULES FOR THIS REQUEST:
1. STRICTLY follow the course methodology
2. Use ONLY formulas from the course
3. ALWAYS complete your answer (never cut off mid-sentence)
4. Address ALL sub-questions (a, b, c) if present
5. Respond ENTIRELY in English"""
        
        user_prompt = f"""CONTEXT:
{context}
{methodology_section}
{formulas_section}
{notation_instructions}

QUESTION:
{question}

COMPLETE ANSWER (based only on CONTEXT, in English):"""
    else:
        system_prompt += """

🎯 RÈGLES ADDITIONNELLES POUR CETTE REQUÊTE :
1. Respecte STRICTEMENT la méthodologie du cours
2. Utilise UNIQUEMENT les formules du cours
3. Complète TOUJOURS ta réponse (jamais de phrase coupée)
4. Traite TOUTES les sous-questions (a, b, c)
5. Réponds ENTIÈREMENT en français"""

        user_prompt = f"""CONTEXT :
{context}
{methodology_section}
{formulas_section}
{notation_instructions}

QUESTION :
{question}

RÉPONSE COMPLÈTE (basée uniquement sur le CONTEXT, en français) :"""

    # ===== GÉNÉRATION CLAUDE =====
    print(f"\n🤖 Génération réponse avec Claude {get_claude_model(user_plan)} en {language.upper()}...")
    
    model = get_claude_model(user_plan)
    
    claude_response = claude_client.messages.create(
        model=model,
        max_tokens=1200,
        temperature=0.1,
        system=system_prompt,
        messages=[
            {"role": "user", "content": user_prompt}
        ]
    )
    
    answer_text = claude_response.content[0].text
    tokens_used = claude_response.usage.input_tokens + claude_response.usage.output_tokens
    
    # Post-traitement des notations mathématiques
    from validation_notations import convert_latex_to_unicode
    
    # Convertir tous les symboles LaTeX en Unicode
    answer_text = convert_latex_to_unicode(answer_text)
    
    # Remplacements supplémentaires pour * (qui n'est pas toujours LaTeX)
    answer_text = answer_text.replace(' * ', ' · ')
    answer_text = answer_text.replace('*', '·')
    
    # Validation
    print("\n✅ Validation réponse...")
    
    validation = validate_answer_completeness(
        question=question,
        answer=answer_text
    )
    
    # Validation des notations
    notation_validation = validate_notation_consistency(answer_text, context_notations)
    
    print(f"   Score de complétude : {validation['score']}/100")
    print(f"   Score de notation : {notation_validation['score']}/100")
    
    if notation_validation['issues']:
        print(f"   ⚠️ Problèmes de notation détectés : {notation_validation['issues_count']}")
        for issue in notation_validation['issues']:
            print(f"      - {issue.get('message','') if isinstance(issue, dict) else str(issue)}")
    
    # Si incomplet, régénérer
    if not validation['is_complete'] and validation['score'] < 50:
        print("\n⚠️ Réponse incomplète - Régénération...")
        
        retry_prompt = user_prompt + f"\n\n⚠️ ATTENTION : Réponse COMPLÈTE requise. Problème : {(validation['issues'][0].get('message','') if isinstance(validation['issues'][0], dict) else str(validation['issues'][0])) if validation['issues'] else 'N/A'}"
        
        claude_response = claude_client.messages.create(
            model=model,
            max_tokens=1500,
            temperature=0.1,
            system=system_prompt,
            messages=[
                {"role": "user", "content": retry_prompt}
            ]
        )
        
        answer_text = claude_response.content[0].text
        tokens_used += claude_response.usage.input_tokens + claude_response.usage.output_tokens
        validation = validate_answer_completeness(question, answer_text)
        
        print(f"   Nouveau score : {validation['score']}/100")
    
    # Confiance
    avg_distance = float(np.mean(distances[0]))
    confidence_rag = 1.0 / (1.0 + avg_distance)
    confidence_final = (confidence_rag * 0.6) + (validation['score'] / 100 * 0.4)
    
    response_time = time.time() - start_time
    
    print(f"\n{'='*60}")
    print(f"✅ RÉPONSE GÉNÉRÉE EN {response_time:.1f}s")
    print(f"🤖 Modèle utilisé : {model}")
    print(f"📊 Tokens utilisés : {tokens_used}")
    print(f"📊 Confiance finale : {confidence_final:.2%}")
    print(f"{'='*60}\n")
    
    result = {
        'question': question,
        'answer': answer_text,
        'sources': [
            {
                'text': c['text'][:200] + '...' if len(c['text']) > 200 else c['text'],
                'page': c['page'],
                'confidence': float(1.0 / (1.0 + c['distance']))
            }
            for c in relevant_chunks[:3]
        ],
        'confidence': confidence_final,
        'response_time': response_time,
        'validation': validation,
        'notation_validation': notation_validation,
        'exercise_type': exercise_type,
        'has_methodology': methodologies_data['has_methodology'],
        'model_used': model,
        'tokens_used': tokens_used,
        'user_plan': user_plan
    }
    
    return result


# ============================================
# MAIN (Pour tests)
# ============================================

if __name__ == "__main__":
    print("🧪 Module RAG chargé avec Claude")
    print(f"✅ Tesseract disponible : {has_tesseract()}")
    print(f"✅ Claude Haiku (free) : {PLAN_MODELS['free']}")
    print(f"✅ Claude Sonnet (payant) : {PLAN_MODELS['pro']}")
